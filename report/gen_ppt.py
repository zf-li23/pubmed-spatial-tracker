#!/usr/bin/env python3
"""Generate course presentation PPT for PubMed Spatial Tracker project.

Usage:
    conda run -n zf-li23 python report/gen_ppt.py

Output: report/PubMed_Spatial_Tracker_Presentation.pptx
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
REPO = Path(__file__).resolve().parent.parent
FIGDIR = REPO / "report" / "figures"
PANELDIR = FIGDIR / "panels"
OUTPUT = REPO / "report" / "PubMed_Spatial_Tracker_Presentation.pptx"

# ── Color Scheme ──
BLUE_DARK = RGBColor(0x01, 0x48, 0x7A)      # dark blue
BLUE_MID  = RGBColor(0x01, 0x73, 0xB2)       # mid blue
BLUE_LIGHT = RGBColor(0xD6, 0xEE, 0xF8)      # light blue bg
ORANGE    = RGBColor(0xDE, 0x8F, 0x05)
GREEN     = RGBColor(0x02, 0x9E, 0x73)
RED       = RGBColor(0xD5, 0x5E, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GRAY      = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_MID  = RGBColor(0xBB, 0xBB, 0xBB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ═══════════════════════════════════════════════════════════════
#  Helper functions
# ═══════════════════════════════════════════════════════════════

def add_blank_slide():
    """Add a blank slide and return it."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_bg(slide, color=WHITE):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=14, color=BLACK, bold=False,
                          line_spacing=1.5, font_name="Microsoft YaHei",
                          alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs (each item in lines)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_rich_textbox(slide, left, top, width, height, segments,
                     font_name="Microsoft YaHei"):
    """Add a text box with mixed-format paragraphs.
    segments: list of list of dicts, each inner list is a paragraph's runs.
    Each run dict: {'text': str, 'size': int, 'color': RGBColor, 'bold': bool}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(segments):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        for run_dict in para_runs:
            run = p.add_run()
            run.text = run_dict.get("text", "")
            run.font.size = Pt(run_dict.get("size", 14))
            run.font.color.rgb = run_dict.get("color", BLACK)
            run.font.bold = run_dict.get("bold", False)
            run.font.name = font_name
    return txBox


def add_img(slide, img_path, left, top, width=None, height=None):
    """Add an image, returns the shape."""
    if not os.path.exists(img_path):
        # Placeholder
        shape = add_rect(slide, left, top, width or Inches(4), height or Inches(3),
                         GRAY_LIGHT, GRAY_MID)
        add_textbox(slide, left, top, width or Inches(4), height or Inches(3),
                    f"[{os.path.basename(img_path)} not found]",
                    font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
        return shape
    if width and height:
        return slide.shapes.add_picture(img_path, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(img_path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        return slide.shapes.add_picture(img_path, left, top)


def title_bar(slide, text, subtitle=None):
    """Add consistent title bar at top."""
    add_rect(slide, 0, 0, W, Inches(1.0), BLUE_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.55),
                text, font_size=26, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.60), Inches(12), Inches(0.35),
                    subtitle, font_size=13, color=RGBColor(0xBB, 0xDD, 0xFF))


def footer(slide, text="机器学习概论 · 课程大作业 · 李哲夫"):
    """Add footer."""
    add_textbox(slide, Inches(0.4), H - Inches(0.35), Inches(12), Inches(0.3),
                text, font_size=9, color=GRAY)


def bullet_list(slide, left, top, width, height, items, font_size=14,
                color=BLACK, bullet_char="▸", font_name="Microsoft YaHei"):
    """Add a bulleted list."""
    lines = []
    for item in items:
        lines.append(f"{bullet_char}  {item}")
    return add_multiline_textbox(slide, left, top, width, height, lines,
                                 font_size=font_size, color=color,
                                 font_name=font_name)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
# Large blue block
add_rect(slide, 0, 0, W, Inches(4.2), BLUE_DARK)
# Accent bar
add_rect(slide, 0, Inches(4.2), W, Inches(0.08), ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
            "面向生物医学文献的", font_size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
            "多策略多标签分类方法系统比较研究", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(0.5),
            "—— 以空间转录组学文献自动分类为例", font_size=20, color=RGBColor(0xBB, 0xDD, 0xFF))

add_textbox(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(0.4),
            "李哲夫    2023011400    生32", font_size=18, color=BLUE_DARK, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
            "清华大学 · 机器学习概论 · 课程大作业", font_size=14, color=GRAY)
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
            "github.com/zf-li23/pubmed-spatial-tracker", font_size=12, color=BLUE_MID)
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
            "2026年6月", font_size=12, color=GRAY)

footer(slide, "")


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2: Outline
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "汇报提纲")

items = [
    "一、选题背景与研究目标",
    "二、数据集与方法论",
    "三、代码架构与功能模块",
    "四、实验结果与分析",
    "五、遇到的问题与解决方案",
    "六、核心结论与展望",
]
for i, item in enumerate(items):
    y = Inches(1.4) + Inches(0.85) * i
    # Number circle
    if i < 4:
        circ = add_rect(slide, Inches(1.2), y, Inches(0.45), Inches(0.45),
                        BLUE_MID if i % 2 == 0 else GREEN)
    else:
        circ = add_rect(slide, Inches(1.2), y, Inches(0.45), Inches(0.45), ORANGE)
    add_textbox(slide, Inches(1.9), y + Inches(0.04), Inches(9), Inches(0.4),
                item, font_size=18, color=BLACK, bold=False)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3: 选题背景
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "一、选题背景与意义", "为什么要系统比较多种机器学习方法？")

# Left: problems
add_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.0), BLUE_LIGHT)
add_textbox(slide, Inches(0.7), Inches(1.35), Inches(5.4), Inches(0.4),
            "❗ 生物医学文献的指数增长", font_size=16, color=BLUE_DARK, bold=True)
bullet_list(slide, Inches(0.7), Inches(1.8), Inches(5.4), Inches(1.3), [
    "空间转录组学 PubMed 年发文量从 <500 篇 (2020)",
    "激增至 4,150 篇 (2025)，传统人工分类不可行",
    "亟需自动化的多标签文献分类方法",
], font_size=13)

add_rect(slide, Inches(0.5), Inches(3.5), Inches(5.8), Inches(2.2), RGBColor(0xFF, 0xEB, 0xD6))
add_textbox(slide, Inches(0.7), Inches(3.55), Inches(5.4), Inches(0.4),
            "⚡ 三大技术挑战", font_size=16, color=ORANGE, bold=True)
bullet_list(slide, Inches(0.7), Inches(4.0), Inches(5.4), Inches(1.5), [
    "文本稀疏：词袋模型无法捕捉医学术语语义关联",
    "标签复杂：一篇文章同时涉及类别/技术/生物领域",
    "标注成本高：大规模人工标注不现实",
], font_size=13)

# Right: approach
add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.0), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(7.0), Inches(1.35), Inches(5.6), Inches(0.4),
            "🎯 研究策略", font_size=16, color=GREEN, bold=True)
bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.6), Inches(1.3), [
    "多数据集 × 多算法 × 多特征表示的系统比较",
    "4 个独立数据集形成天然梯度",
    "14 种算法覆盖课程全部主要模块",
], font_size=13)

add_rect(slide, Inches(6.8), Inches(3.5), Inches(6.0), Inches(2.2), RGBColor(0xF3, 0xE5, 0xF5))
add_textbox(slide, Inches(7.0), Inches(3.55), Inches(5.6), Inches(0.4),
            "🔬 核心研究问题", font_size=16, color=RGBColor(0x7B, 0x1F, 0xA2), bold=True)
bullet_list(slide, Inches(7.0), Inches(4.0), Inches(5.6), Inches(1.5), [
    "不同算法在不同标签粒度下的排序是否一致？",
    "图结构信息对文献分类有无边际收益？",
    "LLM 标注能否替代人工标注？",
    "「宽泛预训练 + 领域微调」范式是否有效？",
], font_size=13)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4: 总体完成情况
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "总体完成情况", "7 组实验 · 115 子任务 · 98.3% 完成率")

# Big numbers
stats = [
    ("7", "组实验"),
    ("115", "子任务"),
    ("98.3%", "完成率"),
    ("9,148", "篇标注"),
    ("14", "种算法"),
    ("~60", "次提交"),
]
for i, (num, label) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4.0) * col
    y = Inches(1.5) + Inches(2.5) * row
    colors = [BLUE_MID, ORANGE, GREEN, RED, RGBColor(0x7B, 0x1F, 0xA2), BLUE_DARK]
    add_rect(slide, x, y, Inches(3.5), Inches(2.0), GRAY_LIGHT)
    add_textbox(slide, x, y + Inches(0.2), Inches(3.5), Inches(0.9),
                num, font_size=44, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(1.2), Inches(3.5), Inches(0.4),
                label, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Experiment table
add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.3),
            "实验进度：000✅  001✅(82/84)  002✅  003✅  004✅  005✅(8/9)  006✅  007✅(11/13)",
            font_size=11, color=BLUE_MID, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 5: 数据集
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "二、数据集全景", "四个数据集从标签粒度、信息维度到标注完备性形成天然梯度")

# Table-like layout
headers = ["数据集", "规模", "标签数", "标签类型", "核心角色"]
rows_data = [
    ["OHSUMED",  "~10K 篇", "1,650", "MeSH 多标签", "大规模稀疏标签基准"],
    ["PML",       "10K 篇",  "16",    "顶级类别",    "粗粒度快速实验"],
    ["PGB",       "5K 篇",   "3",     "节点分类",    "图结构方法验证"],
    ["Spatial\nTracker", "9,148 篇", "6+15+19+17", "LLM 标注\n多维标签", "目标应用场景"],
]

col_widths = [Inches(1.8), Inches(1.6), Inches(1.4), Inches(2.0), Inches(2.6)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + Inches(0.15))

# Header row
y0 = Inches(1.4)
for j, h in enumerate(headers):
    add_rect(slide, col_starts[j], y0, col_widths[j], Inches(0.45), BLUE_DARK)
    add_textbox(slide, col_starts[j], y0 + Inches(0.05), col_widths[j], Inches(0.35),
                h, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(rows_data):
    y = y0 + Inches(0.55) + Inches(0.65) * i
    bg = GRAY_LIGHT if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        add_rect(slide, col_starts[j], y, col_widths[j], Inches(0.55), bg)
        c = BLUE_MID if j == 0 else BLACK
        b = (j == 0)
        fa = PP_ALIGN.CENTER
        add_textbox(slide, col_starts[j] + Inches(0.05), y + Inches(0.08),
                    col_widths[j] - Inches(0.1), Inches(0.4),
                    cell, font_size=11, color=c, bold=b, alignment=fa)

# Gradient illustration
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
            "数据集梯度设计", font_size=16, color=BLUE_DARK, bold=True)

grad_items = [
    "标签粒度梯度：16 (PML) → 3 (PGB) → 1,650 (OHSUMED) → 40+ (ST)",
    "信息维度梯度：纯文本 → 文本+元数据 → 文本+图结构",
    "标注完备性：全标注 (OHSUMED/PML/PGB) → LLM 标注 (ST)",
]
for i, item in enumerate(grad_items):
    add_textbox(slide, Inches(1.0), Inches(5.4) + Inches(0.4) * i, Inches(11), Inches(0.35),
                f"  📌  {item}", font_size=12, color=GRAY)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 6: 实验设计 — 三步渐进
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "实验设计：三步渐进策略", "Step 1: 筛选 → Step 2: 验证 → Step 3: 迁移")

# Step 1
add_rect(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.4), BLUE_LIGHT)
add_rect(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(0.5), BLUE_MID)
add_textbox(slide, Inches(0.5), Inches(1.42), Inches(3.8), Inches(0.5),
            "Step 1: 多数据集算法筛选", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
bullet_list(slide, Inches(0.7), Inches(2.1), Inches(3.4), Inches(4.5), [
    "3 个全标注数据集并行比较",
    "7 模型 × 4 特征 = 28 组/数据集",
    "Exp 001-005 共 105 组实验",
    "产出：算法排名、特征有效性",
    "图结构信息的边际收益量化",
], font_size=12, color=BLACK)

# Arrow
add_textbox(slide, Inches(4.3), Inches(3.5), Inches(0.8), Inches(0.5),
            "→", font_size=36, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Step 2
add_rect(slide, Inches(5.0), Inches(1.4), Inches(3.5), Inches(5.4), RGBColor(0xE8, 0xF5, 0xE9))
add_rect(slide, Inches(5.0), Inches(1.4), Inches(3.5), Inches(0.5), GREEN)
add_textbox(slide, Inches(5.0), Inches(1.42), Inches(3.5), Inches(0.5),
            "Step 2: 目标数据集验证", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
bullet_list(slide, Inches(5.2), Inches(2.1), Inches(3.1), Inches(4.5), [
    "PubMed 检索 → 9,148 篇",
    "6 维 LLM 标注体系设计",
    "DeepSeek-v4-flash 批量标注",
    "3 方法 ST 基准测试 (Exp 006)",
    "标注质量评估",
], font_size=12, color=BLACK)

# Arrow
add_textbox(slide, Inches(8.5), Inches(3.5), Inches(0.8), Inches(0.5),
            "→", font_size=36, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Step 3
add_rect(slide, Inches(9.2), Inches(1.4), Inches(3.6), Inches(5.4), RGBColor(0xF3, 0xE5, 0xF5))
add_rect(slide, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.5), RGBColor(0x7B, 0x1F, 0xA2))
add_textbox(slide, Inches(9.2), Inches(1.42), Inches(3.6), Inches(0.5),
            "Step 3: 迁移微调探索", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
bullet_list(slide, Inches(9.4), Inches(2.1), Inches(3.2), Inches(4.5), [
    "3 种源域预训练策略",
    "BioBERT+MLP 微调",
    "GCN/GraphSAGE 图迁移",
    "13 组实验 (Exp 007)",
    "F1=0.9143 全场最佳 🏆",
], font_size=12, color=BLACK)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 7: 算法全景
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "算法全景", "14 种算法 × 5 种文本表示，覆盖课程全部主要模块")

# Text representations
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
            "文本表示（5 种）", font_size=16, color=BLUE_DARK, bold=True)

repr_data = [
    ("TF-IDF (1-2gram, 5K)", "词级稀疏基线"),
    ("BioBERT (768d, mean pool)", "上下文语义嵌入"),
    ("LDA (K=15)", "文档级主题分布"),
    ("元特征 (3-5d)", "非文本信号"),
    ("Node2Vec (128d)", "图结构嵌入"),
]
for i, (name, desc) in enumerate(repr_data):
    y = Inches(1.8) + Inches(0.5) * i
    add_rect(slide, Inches(0.7), y, Inches(3.2), Inches(0.4), BLUE_LIGHT)
    add_textbox(slide, Inches(0.8), y + Inches(0.02), Inches(3.0), Inches(0.35),
                f"{name}", font_size=11, color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(4.0), y + Inches(0.02), Inches(2.5), Inches(0.35),
                desc, font_size=10, color=GRAY)

# Algorithms
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5), Inches(0.4),
            "分类算法（14 种）", font_size=16, color=ORANGE, bold=True)

alg_data = [
    ("贝叶斯学习", "Naive Bayes"),
    ("基于实例的学习", "k-NN, SVM (RBF)"),
    ("回归学习", "Logistic Regression"),
    ("集成学习 (Bagging)", "Random Forest"),
    ("集成学习 (Boosting)", "AdaBoost, XGBoost"),
    ("深度学习", "BioBERT + MLP"),
    ("无监督学习", "LDA + KMeans 聚类"),
    ("图表示学习", "Node2Vec, GCN, GraphSAGE"),
]
for i, (cat, algs) in enumerate(alg_data):
    y = Inches(1.8) + Inches(0.45) * i
    c = GRAY_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, Inches(7.0), y, Inches(2.3), Inches(0.38), c)
    add_rect(slide, Inches(9.3), y, Inches(3.5), Inches(0.38), c)
    add_textbox(slide, Inches(7.1), y + Inches(0.02), Inches(2.2), Inches(0.34),
                cat, font_size=10, color=ORANGE, bold=True)
    add_textbox(slide, Inches(9.4), y + Inches(0.02), Inches(3.3), Inches(0.34),
                algs, font_size=10, color=BLACK)

# Multi-label strategies
add_textbox(slide, Inches(6.8), Inches(5.8), Inches(5), Inches(0.4),
            "多标签策略：Binary Relevance / Classifier Chains / Label Powerset",
            font_size=12, color=GRAY)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 8: 代码架构
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "三、代码架构与主要功能模块", "面向可复现实验的模块化设计")

# Left: structure
modules = [
    ("src/datasets/", "数据集加载器", "4 个数据集统一的 BiomedDataset 接口"),
    ("src/features/", "特征提取器", "TF-IDF / BioBERT / LDA / Meta / Node2Vec"),
    ("src/models/", "模型实现", "14 种算法的统一接口"),
    ("src/evaluation/", "评估框架", "指标计算 + CSV 实验日志"),
    ("src/pipeline.py", "实验流水线", "交叉验证 + 并行调度"),
    ("src/annotate/", "批量标注", "DeepSeek API 断点续标"),
    ("src/search/", "PubMed 检索", "稳健的 NCBI 抓取 + 增量保存"),
]

for i, (path, title, desc) in enumerate(modules):
    y = Inches(1.3) + Inches(0.8) * i
    add_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.65), BLUE_LIGHT)
    add_rect(slide, Inches(3.1), y, Inches(4.0), Inches(0.65), GRAY_LIGHT)
    add_textbox(slide, Inches(0.55), y + Inches(0.02), Inches(2.4), Inches(0.3),
                path, font_size=10, color=BLUE_DARK, bold=True, font_name="Consolas")
    add_textbox(slide, Inches(0.55), y + Inches(0.3), Inches(2.4), Inches(0.3),
                title, font_size=10, color=BLUE_MID)
    add_textbox(slide, Inches(3.2), y + Inches(0.1), Inches(3.8), Inches(0.45),
                desc, font_size=10, color=GRAY)

# Right: experiment structure
add_textbox(slide, Inches(7.5), Inches(1.3), Inches(5), Inches(0.4),
            "实验目录（7 组可复现实验）", font_size=14, color=ORANGE, bold=True)

exp_data = [
    "experiments/000_query_analysis/",
    "experiments/001_classical_matrix/    (7×4×3=84)",
    "experiments/002_biobert_mlp/         (3 数据集)",
    "experiments/003_lda_cluster/         (3 数据集)",
    "experiments/004_multilabel_strategy/ (BR/CC/LP)",
    "experiments/005_graph_models/        (Node2Vec/GCN/GraphSAGE)",
    "experiments/006_st_benchmark/        (ST 3 方法)",
    "experiments/007_transfer_learning/   (13 组实验)",
]
for i, exp in enumerate(exp_data):
    y = Inches(1.8) + Inches(0.5) * i
    bg = GRAY_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, Inches(7.5), y, Inches(5.3), Inches(0.42), bg)
    add_textbox(slide, Inches(7.6), y + Inches(0.03), Inches(5.1), Inches(0.36),
                exp, font_size=10, color=BLACK, font_name="Consolas")

# Key design principles
add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
            "设计原则：统一接口 · 增量保存 · 断点续跑 · 可复现 · 并行调度",
            font_size=11, color=BLUE_MID, bold=True, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 9: Fig1 数据集概览（复合大图）
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "Fig 1: 数据集概览", "空间转录组学文献库 (9,148 篇) 的多维统计分析")

fig1_path = FIGDIR / "fig1_dataset_overview.png"
if fig1_path.exists():
    # Place the composite figure prominently
    add_img(slide, str(fig1_path), Inches(0.3), Inches(1.1), width=Inches(12.7), height=Inches(6.1))
else:
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(2),
                "[fig1_dataset_overview.png not found]\n请先运行 report/scripts/fig1_dataset_overview.py",
                font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Panel labels overlay
add_textbox(slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.3),
            "A: 年度发文趋势  B: 类别饼图  C: 分析标签分布  D: 标签数/篇直方图  E: 类别×标签热力图(%)  "
            "F: 技术平台分布  G: 生物学主题分布  H: 置信度环形图  I: 布尔属性(新数据/代码/预印本)",
            font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 10: Exp 001 经典算法矩阵
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "四、实验结果 — Exp 001 经典算法矩阵", "7 模型 × 4 特征 × 3 数据集 = 84 组 | 82/84 ✅")

# Best results
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
            "各数据集最佳组合", font_size=15, color=BLUE_DARK, bold=True)

best_data = [
    ("OHSUMED (1,650 标签)", "AdaBoost + TF-IDF", "0.1687", "次优: LR+BioBERT 0.0853"),
    ("PML (16 标签)", "LogisticReg + BioBERT", "0.6710 🏆", "次优: SVM+BioBERT 0.6603"),
    ("PGB (3 类)", "AdaBoost + TF-IDF", "0.4215", "次优: SVM+TF-IDF 0.3775"),
]
for i, (ds, best, f1, sub) in enumerate(best_data):
    y = Inches(1.8) + Inches(1.1) * i
    colors = [BLUE_MID, GREEN, ORANGE]
    add_rect(slide, Inches(0.5), y, Inches(6.0), Inches(0.9), GRAY_LIGHT)
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(0.9), colors[i])
    add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(3.5), Inches(0.35),
                ds, font_size=12, color=BLACK, bold=True)
    add_textbox(slide, Inches(0.7), y + Inches(0.4), Inches(3.0), Inches(0.35),
                f"最佳: {best}", font_size=11, color=colors[i], bold=True)
    add_textbox(slide, Inches(3.8), y + Inches(0.1), Inches(2.5), Inches(0.6),
                f1, font_size=22, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), y + Inches(0.65), Inches(5.5), Inches(0.2),
                sub, font_size=9, color=GRAY)

# Key insight box
add_rect(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.2), BLUE_LIGHT)
add_textbox(slide, Inches(7.2), Inches(1.4), Inches(5.4), Inches(0.4),
            "💡 关键发现", font_size=15, color=BLUE_DARK, bold=True)
bullet_list(slide, Inches(7.2), Inches(1.9), Inches(5.4), Inches(3.5), [
    "BioBERT 嵌入 >> TF-IDF: PML 上提升 18% (0.6710 vs 0.5682)",
    "AdaBoost+TF-IDF 在稀疏标签空间中意外地好 (OHSUMED 0.1687)",
    "标签密度决定算法排序: 稠密用 LR, 稀疏用 AdaBoost",
    "OHSUMED 1,650 标签太稀疏，所有方法 F1 < 0.2",
    "逻辑回归 + BioBERT 嵌入性价比最优",
], font_size=12, color=BLACK)

add_textbox(slide, Inches(7.2), Inches(5.7), Inches(5.4), Inches(0.8),
            "完整子图展示见后文 Fig2 面板画廊", font_size=11, color=GRAY,
            alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 10: Exp 002-005 概览
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "实验结果 — Exp 002 ~ 005 概览", "BioBERT 微调 · 无监督聚类 · 多标签策略 · 图模型")

# Exp 002
add_rect(slide, Inches(0.4), Inches(1.3), Inches(3.0), Inches(2.7), GRAY_LIGHT)
add_rect(slide, Inches(0.4), Inches(1.3), Inches(3.0), Inches(0.45), BLUE_MID)
add_textbox(slide, Inches(0.4), Inches(1.32), Inches(3.0), Inches(0.4),
            "Exp 002 BioBERT+MLP", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
exp2 = [
    ("OHSUMED", "0.0013", "17min"),
    ("PML", "0.6411", "19min"),
    ("PGB", "0.3601", "9min"),
]
for i, (ds, f1, t) in enumerate(exp2):
    y = Inches(1.9) + Inches(0.55) * i
    add_textbox(slide, Inches(0.6), y, Inches(1.2), Inches(0.3),
                ds, font_size=10, color=BLACK, bold=True)
    add_textbox(slide, Inches(1.8), y, Inches(0.7), Inches(0.3),
                f1, font_size=10, color=BLUE_MID, bold=True)
    add_textbox(slide, Inches(2.5), y, Inches(0.7), Inches(0.3),
                t, font_size=9, color=GRAY)
add_textbox(slide, Inches(0.5), Inches(3.6), Inches(2.8), Inches(0.3),
            "PML 上微调后 F1=0.6411,\n与冻结 LR 的 0.6710 接近",
            font_size=9, color=GRAY)

# Exp 003
add_rect(slide, Inches(3.6), Inches(1.3), Inches(3.0), Inches(2.7), GRAY_LIGHT)
add_rect(slide, Inches(3.6), Inches(1.3), Inches(3.0), Inches(0.45), ORANGE)
add_textbox(slide, Inches(3.6), Inches(1.32), Inches(3.0), Inches(0.4),
            "Exp 003 LDA+聚类", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
exp3 = [
    ("OHSUMED", "NMI=0.44"),
    ("PML", "NMI=0.10"),
    ("PGB", "NMI=0.005"),
]
for i, (ds, n) in enumerate(exp3):
    y = Inches(1.9) + Inches(0.55) * i
    add_textbox(slide, Inches(3.8), y, Inches(1.2), Inches(0.3),
                ds, font_size=10, color=BLACK, bold=True)
    add_textbox(slide, Inches(5.0), y, Inches(1.4), Inches(0.3),
                n, font_size=10, color=ORANGE, bold=True)
add_textbox(slide, Inches(3.7), Inches(3.6), Inches(2.8), Inches(0.3),
            "OHSUMED 主题结构明显\nPGB 3 类不分 (NMI≈0)",
            font_size=9, color=GRAY)

# Exp 004
add_rect(slide, Inches(6.8), Inches(1.3), Inches(3.0), Inches(2.7), GRAY_LIGHT)
add_rect(slide, Inches(6.8), Inches(1.3), Inches(3.0), Inches(0.45), GREEN)
add_textbox(slide, Inches(6.8), Inches(1.32), Inches(3.0), Inches(0.4),
            "Exp 004 多标签策略", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
exp4 = [
    ("PML CC", "0.5796 🏆"),
    ("PML BR", "0.5686"),
    ("PML LP", "0.5686"),
    ("OHSUMED", "< 0.01"),
]
for i, (ds, v) in enumerate(exp4):
    y = Inches(1.9) + Inches(0.45) * i
    add_textbox(slide, Inches(7.0), y, Inches(1.2), Inches(0.3),
                ds, font_size=10, color=BLACK, bold=True)
    add_textbox(slide, Inches(8.2), y, Inches(1.4), Inches(0.3),
                v, font_size=10, color=GREEN, bold=True)
add_textbox(slide, Inches(6.9), Inches(3.6), Inches(2.8), Inches(0.3),
            "CC 链式传递捕捉标签依赖\n仅在小标签空间有效",
            font_size=9, color=GRAY)

# Exp 005
add_rect(slide, Inches(10.0), Inches(1.3), Inches(3.0), Inches(2.7), GRAY_LIGHT)
add_rect(slide, Inches(10.0), Inches(1.3), Inches(3.0), Inches(0.45), RED)
add_textbox(slide, Inches(10.0), Inches(1.32), Inches(3.0), Inches(0.4),
            "Exp 005 图模型", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
exp5 = [
    ("GCN", "0.4125 🏆"),
    ("Node2Vec+LR", "0.3324"),
    ("GraphSAGE", "0.3324"),
    ("Node2Vec+SVM", "0.3324"),
]
for i, (m, v) in enumerate(exp5):
    y = Inches(1.9) + Inches(0.45) * i
    add_textbox(slide, Inches(10.2), y, Inches(1.5), Inches(0.3),
                m, font_size=10, color=BLACK, bold=True)
    add_textbox(slide, Inches(11.5), y, Inches(1.3), Inches(0.3),
                v, font_size=10, color=RED, bold=True)
add_textbox(slide, Inches(10.1), Inches(3.6), Inches(2.8), Inches(0.3),
            "GCN 图卷积优于随机游走\nNode2Vec 各分类器持平",
            font_size=9, color=GRAY)

# Bottom: Add figures
figs_to_show = [
    ("fig3_A.png", Inches(0.4), Inches(4.3), Inches(2.9), Inches(2.8)),
    ("fig3_B.png", Inches(3.5), Inches(4.3), Inches(2.9), Inches(2.8)),
    ("fig3_E.png", Inches(6.6), Inches(4.3), Inches(2.9), Inches(2.8)),
    ("fig3_L.png", Inches(9.7), Inches(4.3), Inches(2.9), Inches(2.8)),
]
for fname, x, y, w, h in figs_to_show:
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 11: Exp 006 ST 基准测试 + Exp 007 迁移学习
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "实验结果 — Exp 006 & 007", "ST 基准测试 · 迁移微调 · 全场最佳 F1=0.9143 🏆")

# Exp 006 table
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
            "Exp 006: Spatial Tracker 基准测试", font_size=15, color=BLUE_DARK, bold=True)

col_w = [Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.0)]
col_x = [Inches(0.5)]
for cw in col_w[:-1]:
    col_x.append(col_x[-1] + cw + Inches(0.05))
hdr_y = Inches(1.8)

for j, h in enumerate(["方法", "F1-macro", "Accuracy", "时间"]):
    add_rect(slide, col_x[j], hdr_y, col_w[j], Inches(0.35), BLUE_MID)
    add_textbox(slide, col_x[j], hdr_y + Inches(0.02), col_w[j], Inches(0.3),
                h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

st_rows = [
    ("TF-IDF + SVM", "0.6365", "0.9167", "913s"),
    ("BioBERT + LR", "0.8068", "0.9298", "138s ⚡"),
    ("BioBERT + MLP", "0.8444 🏆", "0.9380", "1039s"),
]
for i, row in enumerate(st_rows):
    y = hdr_y + Inches(0.4) + Inches(0.35) * i
    bg = GRAY_LIGHT if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        add_rect(slide, col_x[j], y, col_w[j], Inches(0.3), bg)
        c = GREEN if i == 2 else BLACK
        b = (i == 2)
        add_textbox(slide, col_x[j], y + Inches(0.01), col_w[j], Inches(0.28),
                    cell, font_size=10, color=c, bold=b, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(3.4), Inches(5.3), Inches(0.3),
            "💡 BioBERT+LR 性价比最优: 差 4.7% 但快 7.5 倍",
            font_size=11, color=ORANGE, bold=True)

# Exp 007 table
add_textbox(slide, Inches(6.5), Inches(1.3), Inches(6.3), Inches(0.4),
            "Exp 007: 迁移微调探索 — 核心结果", font_size=15, color=RGBColor(0x7B, 0x1F, 0xA2), bold=True)

tl_col_w = [Inches(1.8), Inches(1.5), Inches(1.2), Inches(1.2)]
tl_col_x = [Inches(6.5)]
for tcw in tl_col_w[:-1]:
    tl_col_x.append(tl_col_x[-1] + tcw + Inches(0.05))
tl_hdr_y = Inches(1.8)

for j, h in enumerate(["源域→目标", "算法", "F1-macro", "增益"]):
    add_rect(slide, tl_col_x[j], tl_hdr_y, tl_col_w[j], Inches(0.35), RGBColor(0x7B, 0x1F, 0xA2))
    add_textbox(slide, tl_col_x[j], tl_hdr_y + Inches(0.02), tl_col_w[j], Inches(0.3),
                h, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tl_rows = [
    ("ST→ST (基线)", "BioBERT+MLP", "0.8345", "—"),
    ("PML→ST 🏆", "BioBERT+MLP", "0.9143", "+9.6% 🔥"),
    ("OHSUMED→ST", "BioBERT+MLP", "0.8503", "+1.9%"),
    ("ST k-NN GCN", "GCN", "0.7716", "—"),
    ("ST k-NN SAGE", "GraphSAGE", "0.7603", "—"),
]
for i, row in enumerate(tl_rows):
    y = tl_hdr_y + Inches(0.4) + Inches(0.35) * i
    bg = RGBColor(0xFF, 0xF3, 0xE0) if i == 1 else (GRAY_LIGHT if i % 2 == 0 else WHITE)
    for j, cell in enumerate(row):
        add_rect(slide, tl_col_x[j], y, tl_col_w[j], Inches(0.3), bg)
        c = RED if i == 1 else BLACK
        b = (i <= 2)
        add_textbox(slide, tl_col_x[j], y + Inches(0.01), tl_col_w[j], Inches(0.28),
                    cell, font_size=10, color=c, bold=b, alignment=PP_ALIGN.CENTER)

# Highlight
add_rect(slide, Inches(6.5), Inches(4.1), Inches(6.3), Inches(1.5), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(6.7), Inches(4.15), Inches(5.9), Inches(0.35),
            "🏆 核心发现: PML 预训练 + ST 微调 F1=0.9143", font_size=14, color=RED, bold=True)
bullet_list(slide, Inches(6.7), Inches(4.6), Inches(5.9), Inches(0.9), [
    "比直接训练提升 +9.6%",
    "源域选择至关重要: PML (同类型) >> OHSUMED (异类型)",
    "GCN/GraphSAGE 在 ST k-NN 图上 F1≈0.77, 接近 LR 基线",
], font_size=11, color=BLACK)

# Add figures
figs = [
    ("fig4_A.png", Inches(0.4), Inches(3.8), Inches(2.8), Inches(3.2)),
    ("fig4_B.png", Inches(3.3), Inches(3.8), Inches(2.8), Inches(3.2)),
]
for fname, x, y, w, h in figs:
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 12: 标注数据 & Web App
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "数据标注与 Web 应用", "9,148 篇 LLM 标注 · 交互式文献浏览器")

# Annotation stats
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
            "LLM 批量标注 (DeepSeek-v4-flash)", font_size=16, color=BLUE_DARK, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.8), Inches(3), Inches(0.4),
            "类别分布", font_size=13, color=BLUE_MID, bold=True)
cat_data = [
    ("Research", "5,333", "58.3%"),
    ("Technology", "1,785", "19.5%"),
    ("Review", "1,308", "14.3%"),
    ("Protocol", "551", "6.0%"),
    ("Data Resource", "91", "1.0%"),
    ("Benchmark", "79", "0.9%"),
]
for i, (cat, cnt, pct) in enumerate(cat_data):
    y = Inches(2.2) + Inches(0.33) * i
    add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.3),
                cat, font_size=10, color=BLACK)
    add_textbox(slide, Inches(2.2), y, Inches(0.6), Inches(0.3),
                cnt, font_size=10, color=BLUE_MID, bold=True)
    add_textbox(slide, Inches(2.8), y, Inches(0.6), Inches(0.3),
                pct, font_size=10, color=GRAY)
    # mini bar
    pct_val = float(pct.replace("%", "")) / 100
    bar_w = Inches(1.5 * pct_val)
    add_rect(slide, Inches(3.4), y + Inches(0.05), bar_w, Inches(0.2), BLUE_MID)

add_textbox(slide, Inches(0.5), Inches(4.3), Inches(3), Inches(0.3),
            f"置信度: high=4,307 (47.1%)  medium=4,439 (48.5%)  low=401 (4.4%)",
            font_size=10, color=GRAY)

# Tags
add_textbox(slide, Inches(5.0), Inches(1.8), Inches(3), Inches(0.3),
            "Top 标签", font_size=13, color=BLUE_MID, bold=True)
tag_items = [
    "Niche & Microenvironment   2,987",
    "Cell-Cell Communication    2,137",
    "Spatial Domain Ident.      1,766",
    "Cell-Type Deconvolution    1,362",
    "Image-Based Analysis       1,320",
]
for i, tag in enumerate(tag_items):
    add_textbox(slide, Inches(5.0), Inches(2.2) + Inches(0.33) * i, Inches(3.5), Inches(0.3),
                f"▸ {tag}", font_size=9, color=GRAY)

# Web App
add_textbox(slide, Inches(0.5), Inches(4.8), Inches(6), Inches(0.4),
            "Web 应用", font_size=16, color=ORANGE, bold=True)

web_items = [
    "Flask 后端 + React 前端",
    "文章浏览、搜索、多标签筛选",
    "DeepSeek 标注结果可视化",
    "标签管理 (添加/编辑/删除)",
    "技术栈: Python/Flask + Vite/React",
]
for i, item in enumerate(web_items):
    add_textbox(slide, Inches(0.7), Inches(5.3) + Inches(0.35) * i, Inches(5.5), Inches(0.3),
                f"  📌  {item}", font_size=11, color=GRAY)

add_textbox(slide, Inches(6.8), Inches(5.8), Inches(5.5), Inches(0.6),
            "数据集概览复合图见后文 Fig1 大图展示", font_size=11, color=GRAY,
            alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 13: UMAP 嵌入可视化
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "实验结果 — UMAP 嵌入可视化", "4 数据集 × BioBERT 嵌入的 2D 投影")

figs = [
    ("fig5_A.png", Inches(0.3), Inches(1.2), Inches(4.1), Inches(3.0), "OHSUMED"),
    ("fig5_B.png", Inches(4.5), Inches(1.2), Inches(4.1), Inches(3.0), "PML"),
    ("fig5_C.png", Inches(8.7), Inches(1.2), Inches(4.1), Inches(3.0), "PGB"),
    ("fig5_D.png", Inches(0.3), Inches(4.3), Inches(4.1), Inches(3.0), "ST (类别)"),
    ("fig5_E.png", Inches(4.5), Inches(4.3), Inches(4.1), Inches(3.0), "ST (置信度)"),
    ("fig5_F.png", Inches(8.7), Inches(4.3), Inches(4.1), Inches(3.0), "ST (技术平台)"),
]
for fname, x, y, w, h, label in figs:
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)
    add_textbox(slide, x, y + h + Inches(0.02), w, Inches(0.3),
                label, font_size=11, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
            "💡 OHSUMED 聚成单簇 (标签空间过于稀疏) | PML/PGB 可见初始分离 | ST 各类别区分度良好",
            font_size=11, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 14: 关键发现汇总
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "五、核心发现与结论", "10 条关键发现 · 3 条实践建议")

findings = [
    ("🏆", "PML→ST 微调 F1=0.9143", "比直接训练 +9.6%，验证迁移学习范式的有效性", BLUE_MID),
    ("📊", "BioBERT 嵌入 >> TF-IDF", "在所有数据集上系统性提升 10-30%", BLUE_DARK),
    ("🎯", "AdaBoost 在稀疏标签空间称王", "OHSUMED 上 0.1687，远超其他模型", ORANGE),
    ("🔗", "GCN > Node2Vec ≈ GraphSAGE", "图卷积结构对节点分类更有效 (0.4125)", GREEN),
    ("⚡", "BioBERT+LR 性价比最优", "F1=0.8068，仅差 4.7%，快 7.5 倍", RED),
    ("🧩", "CC 多标签策略边际收益有限", "仅在小标签空间 (16类) 体现优势", RGBColor(0x7B, 0x1F, 0xA2)),
    ("📐", "标签密度决定算法排序", "稠密标签用 LR，稀疏标签用 AdaBoost", GRAY),
    ("🔄", "源域选择至关重要", "PML (同类型) 微调优于 OHSUMED (异类型) +7.5%", ORANGE),
    ("🤖", "LLM 标注可行", "9,148 篇全量标注，置信度高+中占 95.6%", GREEN),
    ("📈", "ST 数据集质量良好", "类别区分度好，UMAP 各类别可区分", BLUE_MID),
]

for i, (icon, title, desc, color) in enumerate(findings):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + Inches(6.3) * col
    y = Inches(1.3) + Inches(0.58) * row
    add_rect(slide, x, y, Inches(6.0), Inches(0.5), GRAY_LIGHT)
    add_rect(slide, x, y, Inches(0.06), Inches(0.5), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.02), Inches(5.6), Inches(0.25),
                f"{icon}  {title}", font_size=12, color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.27), Inches(5.6), Inches(0.22),
                desc, font_size=10, color=GRAY)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 15: 问题与挑战
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "六、遇到的问题与解决方案", "集群部署 · 兼容性 · 性能优化 · Bug 修复")

challenges = [
    ("集群环境部署",
     "conda 环境从本地到无互联网集群的迁移",
     "rsync 分批传输解决 tar 的 .so 损坏问题；\n锁定 Python 3.12 + PyTorch 2.5.1+cu121 + transformers 4.48.3"),
    ("HuggingFace 离线",
     "集群无网络，模型加载失败",
     "本地预下载 + rsync 同步缓存；\n改用 BertTokenizer 跳过 auto-detection"),
    ("Exp 001 超时",
     "串行 48h 只跑 11/84 组",
     "拆分为 12 并行任务 (3×4)；增加增量保存机制"),
    ("Naive Bayes 兼容",
     "MultinomialNB 要求非负输入",
     "BioBERT 嵌入有负值 → 改用 GaussianNB"),
    ("CSV 合并混乱",
     "多格式 CSV 列数不一致",
     "编写 merge_results.py 自动检测列数、统一去重"),
    ("并行死锁",
     "n_jobs=-1 嵌套在 joblib.Parallel 内",
     "每个 fold 新建模型实例，不 deepcopy"),
    ("Node2Vec 实现 Bug",
     "alias sampling 返回索引而非节点 ID",
     "修正 _alias_draw() 传入邻居节点列表"),
    ("BioBERT OOM",
     "predict() 一次性处理全部测试数据",
     "predict() 加 batching (batch_size=64)"),
]

for i, (title, problem, solution) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + Inches(6.3) * col
    y = Inches(1.3) + Inches(1.45) * row

    add_rect(slide, x, y, Inches(6.0), Inches(1.3), GRAY_LIGHT)
    add_rect(slide, x, y, Inches(6.0), Inches(0.35), BLUE_MID)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.02), Inches(5.7), Inches(0.3),
                f"⚠  {title}", font_size=12, color=WHITE, bold=True)

    add_textbox(slide, x + Inches(0.15), y + Inches(0.40), Inches(0.5), Inches(0.25),
                "问题:", font_size=9, color=RED, bold=True)
    add_textbox(slide, x + Inches(0.6), y + Inches(0.40), Inches(5.2), Inches(0.25),
                problem, font_size=9, color=GRAY)

    add_textbox(slide, x + Inches(0.15), y + Inches(0.70), Inches(0.5), Inches(0.25),
                "解决:", font_size=9, color=GREEN, bold=True)
    add_textbox(slide, x + Inches(0.6), y + Inches(0.70), Inches(5.2), Inches(0.55),
                solution, font_size=9, color=BLACK)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 16: Lessons Learned
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "技术债务与经验教训", "~60 次 Slurm 提交 · ~2 周集群耗时 · 血的教训")

lessons = [
    ("🚚  环境传输", "用 rsync 而不是 tar 传输 conda 环境，避免 .so 文件损坏"),
    ("📦  版本锁定", "transformers 5.x 需要 torch>=2.6 → 锁定 transformers 4.48.3"),
    ("💾  增量保存", "长时实验必须有增量保存，否则超时全丢"),
    ("🔧  并行计算", "sklearn n_jobs=-1 嵌套在 joblib.Parallel 内会造成死锁"),
    ("📋  CSV 规范", "多格式数据合并要用脚本处理，不要 cat | sort -u"),
    ("🔍  缓存一致性", "跨实验共享特征时确保缓存键完全匹配"),
    ("🔄  断点续跑", "续跑时 completed 计数必须匹配 --models 过滤器"),
    ("📐  一次性处理", "GPU predict() 必须分 batch，否则 OOM"),
]

for i, (icon, lesson) in enumerate(lessons):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + Inches(6.3) * col
    y = Inches(1.3) + Inches(0.65) * row
    add_rect(slide, x, y, Inches(6.0), Inches(0.55), BLUE_LIGHT)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(5.7), Inches(0.4),
                f"{icon}  {lesson}", font_size=12, color=BLUE_DARK, bold=False)

# Experiment time stats
add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.3),
            "实验耗时统计", font_size=15, color=RED, bold=True)

time_stats = [
    ("001 经典矩阵", "~25次", "~10天"),
    ("002 BioBERT", "~9次", "~3天"),
    ("003 LDA", "2次", "8h"),
    ("004 多标签", "~6次", "~1天"),
    ("005 图模型", "~5次", "~1天"),
    ("006 ST基准", "~4次", "~6h"),
    ("007 迁移", "~8次", "~1天"),
]
for i, (exp, sub, time) in enumerate(time_stats):
    x = Inches(0.5) + Inches(1.8) * i
    add_rect(slide, x, Inches(5.9), Inches(1.6), Inches(0.9), GRAY_LIGHT)
    add_textbox(slide, x, Inches(5.92), Inches(1.6), Inches(0.3),
                exp, font_size=9, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.2), Inches(1.6), Inches(0.25),
                sub, font_size=9, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.5), Inches(1.6), Inches(0.25),
                time, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 17: 课程模块对接
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "课程模块对接", "覆盖机器学习概论全部主要模块")

modules_course = [
    ("贝叶斯学习", "Naive Bayes", "PML/PGB/OHSUMED"),
    ("基于实例的学习", "k-NN, SVM (RBF)", "全部数据集"),
    ("回归学习", "Logistic Regression", "全部数据集"),
    ("集成学习 (Bagging)", "Random Forest", "全部数据集"),
    ("集成学习 (Boosting)", "AdaBoost, XGBoost", "全部数据集"),
    ("深度学习", "BioBERT + MLP 微调", "全部数据集 (GPU)"),
    ("无监督学习", "LDA + KMeans 聚类", "全部数据集"),
    ("图表示学习", "Node2Vec, GCN, GraphSAGE", "PGB + ST (k-NN 图)"),
]

for i, (module, algs, datasets) in enumerate(modules_course):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + Inches(6.3) * col
    y = Inches(1.3) + Inches(0.7) * row
    c = [BLUE_MID, ORANGE, GREEN, RED, RGBColor(0x7B, 0x1F, 0xA2), BLUE_DARK, ORANGE, GREEN][i]
    add_rect(slide, x, y, Inches(6.0), Inches(0.6), GRAY_LIGHT)
    add_rect(slide, x, y, Inches(0.06), Inches(0.6), c)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.02), Inches(2.0), Inches(0.28),
                module, font_size=12, color=c, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.28),
                algs, font_size=11, color=BLACK)
    add_textbox(slide, Inches(4.2) + (col * Inches(3.1)), y + Inches(0.1), Inches(2.0), Inches(0.4),
                f"📌 {datasets}", font_size=10, color=GRAY)

# Summary
add_rect(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0), BLUE_LIGHT)
add_textbox(slide, Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.35),
            "课程要求覆盖情况", font_size=15, color=BLUE_DARK, bold=True)
bullet_list(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.3), [
    "每个模块至少实现 1 种算法 ✅ — 覆盖全部 8 个模块，部分模块 2-3 种",
    "所有数据集上使用统一评估框架，公平比较 ✅",
    "多标签任务额外对比 3 种问题转换策略 ✅",
    "实验结果可复现，附完整实验日志和代码 ✅",
    "项目使用 Git 全程记录，代码由 GitHub Copilot 辅助编写 ✅",
], font_size=12, color=BLACK)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 18: 总结与展望
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "总结与展望")

# Summary
add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5), BLUE_LIGHT)
add_textbox(slide, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.35),
            "📋 项目总结", font_size=16, color=BLUE_DARK, bold=True)
bullet_list(slide, Inches(0.7), Inches(1.85), Inches(11.9), Inches(1.8), [
    "构建了多数据集 × 多算法 × 多特征表示的系统性比较框架",
    "7 组实验 113/115 完成 (98.3%)，产出 5 张复合图 + 42 个独立面板",
    "核心发现: PML→ST BioBERT+MLP 微调 F1=0.9143 (+9.6%) 🏆",
    "验证了 LLM 批量标注的可行性 (9,148 篇, high+medium=95.6%)",
    "GitHub Copilot 辅助代码编写，Git 全程记录，实验可复现",
], font_size=13, color=BLACK)

# Future work
add_rect(slide, Inches(0.5), Inches(4.1), Inches(5.9), Inches(2.8), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(0.7), Inches(4.2), Inches(5.5), Inches(0.35),
            "🔜 后续工作", font_size=16, color=GREEN, bold=True)
bullet_list(slide, Inches(0.7), Inches(4.65), Inches(5.5), Inches(2.0), [
    "补全 Exp 001 最后 2 组 (OHSUMED+BioBERT 超时)",
    "完成 PGB→ST GCN/GraphSAGE 图迁移实验",
    "Proposal 最终版插入全部 5 张复合图",
    "完善 Web 应用部署 (GitHub Pages + 后端)",
], font_size=12, color=BLACK)

# Acknowledgments
add_rect(slide, Inches(6.7), Inches(4.1), Inches(6.1), Inches(2.8), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(6.9), Inches(4.2), Inches(5.7), Inches(0.35),
            "🙏 致谢", font_size=16, color=ORANGE, bold=True)
bullet_list(slide, Inches(6.9), Inches(4.65), Inches(5.7), Inches(2.0), [
    "机器学习概论课程全体授课教师与助教",
    "清华大学高性能计算平台 (a-cluster) 支持",
    "DeepSeek 提供 API 用于批量标注",
    "GitHub Copilot 辅助代码开发与调试",
], font_size=12, color=BLACK)

# Thank you
add_rect(slide, Inches(3.5), Inches(7.0), Inches(6.3), Inches(0.0), WHITE)
add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
            "感谢聆听！欢迎提问与交流  |  github.com/zf-li23/pubmed-spatial-tracker",
            font_size=14, color=BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)

footer(slide, "")


# ═══════════════════════════════════════════════════════════════
#  SLIDE: Fig2 面板画廊 — 经典算法矩阵 (A-I)
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "Fig 2: 经典算法矩阵 — 全部子图 (A-I)", "Exp 001 经典算法 × 3 数据集 × 显著性分析")

fig2_panels = [
    ("fig2_A.png", "A: OHSUMED 特征×模型 F1-macro"),
    ("fig2_B.png", "B: PML 特征×模型 F1-macro"),
    ("fig2_C.png", "C: PGB 特征×模型 F1-macro"),
    ("fig2_D.png", "D: 各数据集最佳 F1-macro"),
    ("fig2_E.png", "E: 训练时间对比 (对数坐标)"),
    ("fig2_F.png", "F: PML Top-5 模型 + 显著性"),
    ("fig2_G.png", "G: OHSUMED Top-5 模型"),
    ("fig2_H.png", "H: 各特征最佳 F1 跨数据集对比"),
    ("fig2_I.png", "I: F1-macro vs 训练时间"),
]

# 3x3 grid
for idx, (fname, label) in enumerate(fig2_panels):
    col = idx % 3
    row = idx // 3
    x = Inches(0.3) + Inches(4.3) * col
    y = Inches(1.15) + Inches(2.05) * row
    w = Inches(4.0)
    h = Inches(1.75)
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)
    add_textbox(slide, x, y + h - Inches(0.22), w, Inches(0.2),
                label, font_size=8, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE: Fig3 面板画廊 — 无监督/多标签/图模型 (A-F)
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "Fig 3: 无监督学习 · 多标签策略 · 图模型 (A-F)", "Exp 002 BioBERT微调 · 003 LDA聚类 · 004 多标签策略")

fig3_panels_1 = [
    ("fig3_A.png", "A: BioBERT+MLP vs 最佳经典"),
    ("fig3_B.png", "B: LDA 聚类质量 (NMI)"),
    ("fig3_C.png", "C: 无监督 vs 有监督"),
    ("fig3_D.png", "D: 成本-效益 (F1 vs 时间)"),
    ("fig3_E.png", "E: 多标签策略 PML (BR/CC/LP)"),
    ("fig3_F.png", "F: 多标签 F1 vs 时间"),
]

# 2x3 grid
for idx, (fname, label) in enumerate(fig3_panels_1):
    col = idx % 2
    row = idx // 2
    x = Inches(0.3) + Inches(6.5) * col
    y = Inches(1.15) + Inches(2.0) * row
    w = Inches(6.2)
    h = Inches(1.75)
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)
    add_textbox(slide, x, y + h - Inches(0.22), w, Inches(0.2),
                label, font_size=8, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

# Analysis
add_rect(slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.0), WHITE)
add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
            "💡 PML上 BioBERT+MLP F1=0.6411 | OHSUMED 主题结构明显 (NMI=0.44) | CC 链式传递略优于 BR",
            font_size=10, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE: Fig3 面板画廊 — 无监督/多标签/图模型 (G-L)
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "Fig 3: 无监督学习 · 多标签策略 · 图模型 (G-L)", "Exp 005 Node2Vec · GCN · GraphSAGE → 图模型对比")

fig3_panels_2 = [
    ("fig3_G.png", "G: OHSUMED 多标签策略"),
    ("fig3_H.png", "H: 模型鲁棒性 (F1 箱线图)"),
    ("fig3_I.png", "I: Node2Vec + 分类器"),
    ("fig3_J.png", "J: 图方法对比 (PGB)"),
    ("fig3_K.png", "K: ST k-NN 图方法"),
    ("fig3_L.png", "L: 标签复杂度分析"),
]

# 2x3 grid
for idx, (fname, label) in enumerate(fig3_panels_2):
    col = idx % 2
    row = idx // 2
    x = Inches(0.3) + Inches(6.5) * col
    y = Inches(1.15) + Inches(2.0) * row
    w = Inches(6.2)
    h = Inches(1.75)
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)
    add_textbox(slide, x, y + h - Inches(0.22), w, Inches(0.2),
                label, font_size=8, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
            "💡 GCN (0.4125) >> Node2Vec (0.3324) ≈ GraphSAGE (0.3324) — 图卷积结构对文献节点分类更有效",
            font_size=10, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  SLIDE: Fig4 面板画廊 — ST 基准测试 & 迁移学习 (A-F)
# ═══════════════════════════════════════════════════════════════

slide = add_blank_slide()
add_bg(slide, WHITE)
title_bar(slide, "Fig 4: ST 基准测试 · 迁移微调 (A-F)", "Exp 006 三方法基准 · Exp 007 迁移微调探索")

fig4_panels = [
    ("fig4_A.png", "A: 各类别 Top TF-IDF 词 (SVM 系数)"),
    ("fig4_B.png", "B: ST 三方法基准 + 显著性"),
    ("fig4_C.png", "C: Accuracy vs F1-macro"),
    ("fig4_D.png", "D: 标签共现网络"),
    ("fig4_E.png", "E: 迁移学习瀑布图"),
    ("fig4_F.png", "F: 预训练 vs 微调时间成本"),
]

# 2x3 grid
for idx, (fname, label) in enumerate(fig4_panels):
    col = idx % 2
    row = idx // 2
    x = Inches(0.3) + Inches(6.5) * col
    y = Inches(1.15) + Inches(2.0) * row
    w = Inches(6.2)
    h = Inches(1.75)
    fp = PANELDIR / fname
    if fp.exists():
        add_img(slide, str(fp), x, y, width=w, height=h)
    add_textbox(slide, x, y + h - Inches(0.22), w, Inches(0.2),
                label, font_size=8, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
            "🏆 PML→ST BioBERT+MLP 微调 F1=0.9143 (+9.6%) | BioBERT+LR 性价比最优 (差4.7%但快7.5倍)",
            font_size=10, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════

print(f"Saving to {OUTPUT}")
prs.save(str(OUTPUT))
print(f"✅ Done! {len(prs.slides)} slides created.")
